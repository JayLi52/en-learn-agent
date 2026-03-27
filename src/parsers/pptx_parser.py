"""PPTX 文档解析器"""
from typing import List
from pathlib import Path
from .base import BaseParser, Document


class PPTXParser(BaseParser):
    """解析 PPTX 文档"""

    def parse(self) -> List[Document]:
        self.validate_file()
        documents = []

        from pptx import Presentation

        prs = Presentation(self.file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []

            # 提取幻灯片标题
            if slide.shapes.title:
                texts.append(f"## {slide.shapes.title.text}")

            # 提取所有文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        texts.append(shape.text)

            if texts:
                documents.append(Document(
                    content="\n\n".join(texts),
                    source=str(self.file_path),
                    page=slide_num,
                    metadata={
                        "file_type": "pptx",
                        "slide": slide_num,
                    }
                ))

        return documents